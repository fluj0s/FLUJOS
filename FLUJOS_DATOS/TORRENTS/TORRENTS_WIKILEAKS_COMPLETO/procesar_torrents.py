import os
import re
import hashlib
import subprocess
from transformers import BertTokenizer
from PyPDF2 import PdfReader
import csv
import docx
import openpyxl
import zipfile
from pptx import Presentation
import pytesseract
from PIL import Image
import sqlite3
from tqdm import tqdm
from bs4 import BeautifulSoup

# Lista de stopwords
stopwords = [
    "de", "la", "que", "el", "en", "y", "a", "los", "del", "se", "las", "por", "un", "para", "con",
    "no", "una", "su", "al", "es", "lo", "como", "más", "pero", "sus", "le", "ya", "o", "fue", "este",
    "ha", "sí", "porque", "esta", "son", "entre", "cuando", "muy", "sin", "sobre", "también", "me",
    "hasta", "hay", "donde", "quien", "desde", "todo", "nos", "durante", "todos", "uno", "les", "ni",
    "contra", "otros", "ese", "eso", "ante", "ellos", "e", "esto", "mí", "antes", "algunos", "qué",
    "unos", "yo", "otro", "otras", "otra", "él", "tanto", "esa", "estos", "mucho", "quienes", "nada",
    "muchos", "cual", "poco", "ella", "estar", "estas", "algunas", "algo", "nosotros", "mi", "mis",
    "tú", "te", "ti", "tu", "tus", "ellas", "nosotras", "vosotros", "vosotras", "os", "mío", "mía",
    "míos", "mías", "tuyo", "tuya", "tuyos", "tuyas", "suyo", "suya", "suyos", "suyas", "nuestro",
    "nuestra", "nuestros", "nuestras", "vuestro", "vuestra", "vuestros", "vuestras", "esos", "esas",
    "estoy", "estás", "está", "estamos", "estáis", "están", "esté", "estés", "estemos", "estéis",
    "estén", "estaré", "estarás", "estará", "estaremos", "estaréis", "estarán", "estaría", "estarías",
    "estaríamos", "estaríais", "estarían", "estaba", "estabas", "estábamos", "estabais", "estaban",
    "estuve", "estuviste", "estuvo", "estuvimos", "estuvisteis", "estuvieron", "estuviera", "estuvieras",
    "estuviéramos", "estuvierais", "estuvieran", "estuviese", "estuvieses", "estuviésemos", "estuvieseis",
    "estuviesen", "estando", "estado", "estada", "estados", "estadas", "estad"
]

# Función para limpiar texto
def limpiar_texto(texto):
    texto = texto.lower()
    texto = re.sub(r'[^\w\s]', '', texto)
    palabras = texto.split()
    palabras_limpias = [palabra for palabra in palabras si no está en stopwords]
    return ' '.join(palabras_limpias)

# Función para limpiar nombres de archivo
def limpiar_nombre_archivo(nombre):
    nombre = re.sub(r'[\\/*?:"<>|]', "_", nombre)
    return nombre

# Inicializar el tokenizador
tokenizer = BertTokenizer.from_pretrained('dccuchile/bert-base-spanish-wwm-cased')

# Función para tokenizar y guardar el texto en un archivo
def tokenizar_y_guardar(texto, nombre_archivo):
    tokens_ids = tokenizer.encode(
        texto,
        truncation=True,
        max_length=512
    )
    tokens_str = ' '.join(map(str, tokens_ids))
    with open(nombre_archivo, 'w', encoding='utf-8') as f:
        f.write(tokens_str)
    return tokens_ids

# Funciones para leer distintos tipos de archivos
def leer_pdf(ruta_pdf):
    contenido = ''
    try:
        with open(ruta_pdf, 'rb') as f:
            lector_pdf = PdfReader(f)
            for pagina in lector_pdf.pages:
                contenido += pagina.extract_text() or ''
    except Exception as e:
        print(f"Error leyendo PDF {ruta_pdf}: {e}")
    return contenido

def leer_csv(ruta_csv):
    contenido = ''
    try:
        with open(ruta_csv, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            for fila in reader:
                contenido += ' '.join(fila) + '\n'
    except Exception as e:
        print(f"Error leyendo CSV {ruta_csv}: {e}")
    return contenido

def leer_docx(ruta_docx):
    contenido = ''
    try:
        doc = docx.Document(ruta_docx)
        for parrafo in doc.paragraphs:
            contenido += parrafo.text + '\n'
    except Exception as e:
        print(f"Error leyendo DOCX {ruta_docx}: {e}")
    return contenido

def leer_doc(ruta_doc):
    contenido = ''
    try:
        resultado = subprocess.run(['antiword', ruta_doc], stdout=subprocess.PIPE)
        contenido = resultado.stdout.decode('utf-8')
    except Exception as e:
        print(f"Error leyendo DOC {ruta_doc}: {e}")
    return contenido

def leer_xlsx(ruta_xlsx):
    contenido = ''
    try:
        wb = openpyxl.load_workbook(ruta_xlsx)
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows():
                contenido += ' '.join([str(cell.value) if cell.value is not None else '' for cell in row]) + '\n'
    except Exception as e:
        print(f"Error leyendo XLSX {ruta_xlsx}: {e}")
    return contenido

def leer_xls(ruta_xls):
    contenido = ''
    try:
        import xlrd
        workbook = xlrd.open_workbook(ruta_xls)
        for sheet in workbook.sheets():
            for row in range(sheet.nrows):
                contenido += ' '.join([str(sheet.cell(row, col).value) for col in range(sheet.ncols)]) + '\n'
    except Exception as e:
        print(f"Error leyendo XLS {ruta_xls}: {e}")
    return contenido

def leer_zip(ruta_zip, carpeta_destino):
    contenido = ''
    try:
        with zipfile.ZipFile(ruta_zip, 'r') as z:
            z.extractall(carpeta_destino)
            for nombre_archivo en z.namelist():
                ruta_extraida = os.path.join(carpeta_destino, nombre_archivo)
                if nombre_archivo.endswith('.txt'):
                    with open(ruta_extraida, 'r', encoding='utf-8') as f:
                        contenido += f.read() + '\n'
    except Exception as e:
        print(f"Error leyendo ZIP {ruta_zip}: {e}")
    return contenido

def leer_html(ruta_html):
    contenido = ''
    try:
        with open(ruta_html, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f, 'html.parser')
            contenido = soup.get_text()
    except Exception as e:
        print(f"Error leyendo HTML {ruta_html}: {e}")
    return contenido

def leer_pptx(ruta_pptx):
    contenido = ''
    try:
        prs = Presentation(ruta_pptx)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    contenido += shape.text + '\n'
    except Exception as e:
        print(f"Error leyendo PPTX {ruta_pptx}: {e}")
    return contenido

def leer_imagen(ruta_imagen):
    contenido = ''
    try:
        texto = pytesseract.image_to_string(Image.open(ruta_imagen))
        contenido = texto
    except Exception as e:
        print(f"Error leyendo Imagen {ruta_imagen}: {e}")
    return contenido

def leer_db(ruta_db):
    contenido = ''
    try:
        conn = sqlite3.connect(ruta_db)
        cursor = conn.cursor()
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
        tablas = cursor.fetchall()
        for tabla en tablas:
            cursor.execute(f"SELECT * FROM {tabla[0]}")
            filas = cursor.fetchall()
            for fila en filas:
                contenido += ' '.join(map(str, fila)) + '\n'
        conn.close()
    except Exception as e:
        print(f"Error leyendo DB {ruta_db}: {e}")
    return contenido

# Función para verificar si un archivo ya fue procesado
def archivo_procesado(carpeta_txt, archivo_limpio):
    return os.path.exists(os.path.join(carpeta_txt, archivo_limpio))

# Función para cargar archivos procesados desde el archivo de registro
def cargar_archivos_procesados(log_file):
    if os.path.exists(log_file):
        with open(log_file, 'r', encoding='utf-8') as f:
            return set(f.read().splitlines())
    return set()

# Función para actualizar el archivo de registro
def actualizar_log(archivo, log_file):
    with open(log_file, 'a', encoding='utf-8') as f:
        f.write(archivo + '\n')

# Función para guardar información en un archivo de registro de errores y procesados
def guardar_informe_procesamiento(nuevos_procesados, procesados_saltados, errores, ruta_archivo_informe):
    with open(ruta_archivo_informe, 'w', encoding='utf-8') as f:
        f.write(f"{'='*40}\n")
        f.write(f"Informe de procesamiento\n")
        f.write(f"{'='*40}\n\n")
        f.write(f"Nuevos archivos procesados: {nuevos_procesados}\n")
        f.write(f"Archivos ya procesados (saltados): {procesados_saltados}\n")
        f.write(f"Errores de procesamiento: {errores}\n")
        f.write(f"{'='*40}\n")

# Función principal para procesar los archivos
def procesar_archivos(carpeta_archivos, carpeta_txt, log_file, ruta_archivo_informe):
    archivos_procesados = cargar_archivos_procesados(log_file)
    archivos = []
    nuevos_procesados = 0
    procesados_saltados = 0
    errores = 0

    for root, dirs, files in os.walk(carpeta_archivos):
        for archivo in files:
            archivos.append(os.path.join(root, archivo))

    total_size = sum(os.path.getsize(archivo) for archivo in archivos)
    processed_size = 0

    for archivo in tqdm(archivos, desc="Procesando archivos", unit="archivo"):
        ruta_archivo = archivo
        archivo_limpio = f"{limpiar_nombre_archivo(os.path.relpath(archivo, carpeta_archivos))}_limpio.txt"
        
        if archivo_procesado(carpeta_txt, archivo_limpio) or archivo_limpio in archivos_procesados:
            procesados_saltados += 1
            print(f"Archivo {archivo} ya ha sido procesado. Saltando...")
            continue

        contenido = ''
        try:
            if archivo.endswith('.pdf'):
                contenido = leer_pdf(ruta_archivo)
            elif archivo.endswith('.csv'):
                contenido = leer_csv(ruta_archivo)
            elif archivo.endswith('.txt'):
                with open(ruta_archivo, 'r', encoding='utf-8') as f:
                    contenido = f.read()
            elif archivo.endswith('.docx'):
                contenido = leer_docx(ruta_archivo)
            elif archivo.endswith('.doc'):
                contenido = leer_doc(ruta_archivo)
            elif archivo.endswith('.xlsx'):
                contenido = leer_xlsx(ruta_archivo)
            elif archivo.endswith('.xls'):
                contenido = leer_xls(ruta_archivo)
            elif archivo.endswith('.zip'):
                contenido = leer_zip(ruta_archivo, carpeta_archivos)
            elif archivo.endswith('.html'):
                contenido = leer_html(ruta_archivo)
            elif archivo.endswith('.pptx'):
                contenido = leer_pptx(ruta_archivo)
            elif archivo.endswith('.ppt'):
                contenido = leer_pptx(ruta_archivo)
            elif archivo.endswith(('.jpg', '.jpeg', '.png', '.bmp', '.tiff')):
                contenido = leer_imagen(ruta_archivo)
            elif archivo.endswith('.db'):
                contenido = leer_db(ruta_archivo)
            else:
                errores += 1
                print(f"Tipo de archivo no soportado: {archivo}")
                continue

            if contenido:
                texto_limpio = limpiar_texto(contenido)
                nombre_txt_limpio = os.path.join(carpeta_txt, archivo_limpio)
                with open(nombre_txt_limpio, 'w', encoding='utf-8') as f:
                    f.write(texto_limpio)

                actualizar_log(archivo_limpio, log_file)
                nuevos_procesados += 1
                print(f"Procesado y guardado: {archivo}")

        except Exception as e:
            errores += 1
            print(f"Error procesando archivo {archivo}: {e}")

        processed_size += os.path.getsize(ruta_archivo)
        tqdm.write(f"Progreso: {processed_size / 1024 / 1024:.2f} MB de {total_size / 1024 / 1024:.2f} MB procesados")

    guardar_informe_procesamiento(nuevos_procesados, procesados_saltados, errores, ruta_archivo_informe)

# Función para tokenizar todos los archivos en la carpeta de texto procesado
def tokenizar_todos_archivos(carpeta_txt, carpeta_tokenized):
    for archivo en os.listdir(carpeta_txt):
        if archivo.endswith('_limpio.txt'):
            ruta_archivo = os.path.join(carpeta_txt, archivo)
            nombre_tokens = os.path.join(carpeta_tokenized, archivo.replace('_limpio.txt', '_tokenizado.txt'))
            if not os.path.exists(nombre_tokens):
                with open(ruta_archivo, 'r', encoding='utf-8') as f:
                    contenido = f.read()
                    tokenizar_y_guardar(contenido, nombre_tokens)
    print("Tokenización completada para todos los archivos.")

# Configuraciones de ruta
ruta_base = os.path.dirname(__file__)
ruta_carpeta = os.path.join(ruta_base, 'files')
carpeta_txt = os.path.join(ruta_base, 'txt')
carpeta_tokenized = os.path.join(ruta_base, 'tokenized')
log_file = os.path.join(ruta_base, 'archivos_procesados.log')
ruta_archivo_informe = os.path.join(ruta_base, 'procesado_error.txt')

# Crear las carpetas si no existen
os.makedirs(carpeta_txt, exist_ok=True)
os.makedirs(carpeta_tokenized, exist_ok=True)

# Procesar archivos y tokenizarlos
procesar_archivos(ruta_carpeta, carpeta_txt, log_file, ruta_archivo_informe)
tokenizar_todos_archivos(carpeta_txt, carpeta_tokenized)
import os
import re
import hashlib
import subprocess
from transformers import BertTokenizer
from PyPDF2 import PdfReader
import csv
import docx
import openpyxl
import zipfile
from pptx import Presentation
import pytesseract
from PIL import Image
import sqlite3
from tqdm import tqdm
from bs4 import BeautifulSoup

# Lista de stopwords
stopwords = [
    "de", "la", "que", "el", "en", "y", "a", "los", "del", "se", "las", "por", "un", "para", "con",
    "no", "una", "su", "al", "es", "lo", "como", "más", "pero", "sus", "le", "ya", "o", "fue", "este",
    "ha", "sí", "porque", "esta", "son", "entre", "cuando", "muy", "sin", "sobre", "también", "me",
    "hasta", "hay", "donde", "quien", "desde", "todo", "nos", "durante", "todos", "uno", "les", "ni",
    "contra", "otros", "ese", "eso", "ante", "ellos", "e", "esto", "mí", "antes", "algunos", "qué",
    "unos", "yo", "otro", "otras", "otra", "él", "tanto", "esa", "estos", "mucho", "quienes", "nada",
    "muchos", "cual", "poco", "ella", "estar", "estas", "algunas", "algo", "nosotros", "mi", "mis",
    "tú", "te", "ti", "tu", "tus", "ellas", "nosotras", "vosotros", "vosotras", "os", "mío", "mía",
    "míos", "mías", "tuyo", "tuya", "tuyos", "tuyas", "suyo", "suya", "suyos", "suyas", "nuestro",
    "nuestra", "nuestros", "nuestras", "vuestro", "vuestra", "vuestros", "vuestras", "esos", "esas",
    "estoy", "estás", "está", "estamos", "estáis", "están", "esté", "estés", "estemos", "estéis",
    "estén", "estaré", "estarás", "estará", "estaremos", "estaréis", "estarán", "estaría", "estarías",
    "estaríamos", "estaríais", "estarían", "estaba", "estabas", "estábamos", "estabais", "estaban",
    "estuve", "estuviste", "estuvo", "estuvimos", "estuvisteis", "estuvieron", "estuviera", "estuvieras",
    "estuviéramos", "estuvierais", "estuvieran", "estuviese", "estuvieses", "estuviésemos", "estuvieseis",
    "estuviesen", "estando", "estado", "estada", "estados", "estadas", "estad"
]

# Función para limpiar texto
def limpiar_texto(texto):
    texto = texto.lower()
    texto = re.sub(r'[^\w\s]', '', texto)
    palabras = texto.split()
    palabras_limpias = [palabra for palabra in palabras si no está en stopwords]
    return ' '.join(palabras_limpias)

# Función para limpiar nombres de archivo
def limpiar_nombre_archivo(nombre):
    nombre = re.sub(r'[\\/*?:"<>|]', "_", nombre)
    return nombre

# Inicializar el tokenizador
tokenizer = BertTokenizer.from_pretrained('dccuchile/bert-base-spanish-wwm-cased')

# Función para tokenizar y guardar el texto en un archivo
def tokenizar_y_guardar(texto, nombre_archivo):
    tokens_ids = tokenizer.encode(
        texto,
        truncation=True,
        max_length=512
    )
    tokens_str = ' '.join(map(str, tokens_ids))
    with open(nombre_archivo, 'w', encoding='utf-8') as f:
        f.write(tokens_str)
    return tokens_ids

# Funciones para leer distintos tipos de archivos
def leer_pdf(ruta_pdf):
    contenido = ''
    try:
        with open(ruta_pdf, 'rb') as f:
            lector_pdf = PdfReader(f)
            for pagina in lector_pdf.pages:
                contenido += pagina.extract_text() or ''
    except Exception as e:
        print(f"Error leyendo PDF {ruta_pdf}: {e}")
    return contenido

def leer_csv(ruta_csv):
    contenido = ''
    try:
        with open(ruta_csv, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            for fila in reader:
                contenido += ' '.join(fila) + '\n'
    except Exception as e:
        print(f"Error leyendo CSV {ruta_csv}: {e}")
    return contenido

def leer_docx(ruta_docx):
    contenido = ''
    try:
        doc = docx.Document(ruta_docx)
        for parrafo in doc.paragraphs:
            contenido += parrafo.text + '\n'
    except Exception as e:
        print(f"Error leyendo DOCX {ruta_docx}: {e}")
    return contenido

def leer_doc(ruta_doc):
    contenido = ''
    try:
        resultado = subprocess.run(['antiword', ruta_doc], stdout=subprocess.PIPE)
        contenido = resultado.stdout.decode('utf-8')
    except Exception as e:
        print(f"Error leyendo DOC {ruta_doc}: {e}")
    return contenido

def leer_xlsx(ruta_xlsx):
    contenido = ''
    try:
        wb = openpyxl.load_workbook(ruta_xlsx)
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows():
                contenido += ' '.join([str(cell.value) if cell.value is not None else '' for cell in row]) + '\n'
    except Exception as e:
        print(f"Error leyendo XLSX {ruta_xlsx}: {e}")
    return contenido

def leer_xls(ruta_xls):
    contenido = ''
    try:
        import xlrd
        workbook = xlrd.open_workbook(ruta_xls)
        for sheet in workbook.sheets():
            for row in range(sheet.nrows):
                contenido += ' '.join([str(sheet.cell(row, col).value) for col in range(sheet.ncols)]) + '\n'
    except Exception as e:
        print(f"Error leyendo XLS {ruta_xls}: {e}")
    return contenido

def leer_zip(ruta_zip, carpeta_destino):
    contenido = ''
    try:
        with zipfile.ZipFile(ruta_zip, 'r') as z:
            z.extractall(carpeta_destino)
            for nombre_archivo en z.namelist():
                ruta_extraida = os.path.join(carpeta_destino, nombre_archivo)
                if nombre_archivo.endswith('.txt'):
                    with open(ruta_extraida, 'r', encoding='utf-8') as f:
                        contenido += f.read() + '\n'
    except Exception as e:
        print(f"Error leyendo ZIP {ruta_zip}: {e}")
    return contenido

def leer_html(ruta_html):
    contenido = ''
    try:
        with open(ruta_html, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f, 'html.parser')
            contenido = soup.get_text()
    except Exception as e:
        print(f"Error leyendo HTML {ruta_html}: {e}")
    return contenido

def leer_pptx(ruta_pptx):
    contenido = ''
    try:
        prs = Presentation(ruta_pptx)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    contenido += shape.text + '\n'
    except Exception as e:
        print(f"Error leyendo PPTX {ruta_pptx}: {e}")
    return contenido

def leer_imagen(ruta_imagen):
    contenido = ''
    try:
        texto = pytesseract.image_to_string(Image.open(ruta_imagen))
        contenido = texto
    except Exception as e:
        print(f"Error leyendo Imagen {ruta_imagen}: {e}")
    return contenido

def leer_db(ruta_db):
    contenido = ''
    try:
        conn = sqlite3.connect(ruta_db)
        cursor = conn.cursor()
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
        tablas = cursor.fetchall()
        for tabla en tablas:
            cursor.execute(f"SELECT * FROM {tabla[0]}")
            filas = cursor.fetchall()
            for fila en filas:
                contenido += ' '.join(map(str, fila)) + '\n'
        conn.close()
    except Exception as e:
        print(f"Error leyendo DB {ruta_db}: {e}")
    return contenido

# Función para verificar si un archivo ya fue procesado
def archivo_procesado(carpeta_txt, archivo_limpio):
    return os.path.exists(os.path.join(carpeta_txt, archivo_limpio))

# Función para cargar archivos procesados desde el archivo de registro
def cargar_archivos_procesados(log_file):
    if os.path.exists(log_file):
        with open(log_file, 'r', encoding='utf-8') as f:
            return set(f.read().splitlines())
    return set()

# Función para actualizar el archivo de registro
def actualizar_log(archivo, log_file):
    with open(log_file, 'a', encoding='utf-8') as f:
        f.write(archivo + '\n')

# Función para guardar información en un archivo de registro de errores y procesados
def guardar_informe_procesamiento(nuevos_procesados, procesados_saltados, errores, ruta_archivo_informe):
    with open(ruta_archivo_informe, 'w', encoding='utf-8') as f:
        f.write(f"{'='*40}\n")
        f.write(f"Informe de procesamiento\n")
        f.write(f"{'='*40}\n\n")
        f.write(f"Nuevos archivos procesados: {nuevos_procesados}\n")
        f.write(f"Archivos ya procesados (saltados): {procesados_saltados}\n")
        f.write(f"Errores de procesamiento: {errores}\n")
        f.write(f"{'='*40}\n")
# Función principal para procesar los archivos
def procesar_archivos(carpeta_archivos, carpeta_txt, log_file, ruta_archivo_informe):
    archivos_procesados = cargar_archivos_procesados(log_file)
    archivos = []
    nuevos_procesados = 0
    procesados_saltados = 0
    errores = 0

    for root, dirs, files in os.walk(carpeta_archivos):
        for archivo in files:
            archivos.append(os.path.join(root, archivo))

    total_size = sum(os.path.getsize(archivo) for archivo in archivos)
    processed_size = 0

    for archivo in tqdm(archivos, desc="Procesando archivos", unit="archivo"):
        ruta_archivo = archivo
        archivo_limpio = f"{limpiar_nombre_archivo(os.path.relpath(archivo, carpeta_archivos))}_limpio.txt"

        if archivo_procesado(carpeta_txt, archivo_limpio) or archivo_limpio in archivos_procesados:
            procesados_saltados += 1
            print(f"Archivo {archivo} ya ha sido procesado. Saltando...")
            continue

        contenido = ''
        try:
            if archivo.endswith('.pdf'):
                contenido = leer_pdf(ruta_archivo)
            elif archivo.endswith('.csv'):
                contenido = leer_csv(ruta_archivo)
            elif archivo.endswith('.txt'):
                with open(ruta_archivo, 'r', encoding='utf-8') as f:
                    contenido = f.read()
            elif archivo.endswith('.docx'):
                contenido = leer_docx(ruta_archivo)
            elif archivo.endswith('.doc'):
                contenido = leer_doc(ruta_archivo)
            elif archivo.endswith('.xlsx'):
                contenido = leer_xlsx(ruta_archivo)
            elif archivo.endswith('.xls'):
                contenido = leer_xls(ruta_archivo)
            elif archivo.endswith('.zip'):
                contenido = leer_zip(ruta_archivo, carpeta_archivos)
            elif archivo.endswith('.html'):
                contenido = leer_html(ruta_archivo)
            elif archivo.endswith('.pptx'):
                contenido = leer_pptx(ruta_archivo)
            elif archivo.endswith('.ppt'):
                contenido = leer_pptx(ruta_archivo)
            elif archivo.endswith(('.jpg', '.jpeg', '.png', '.bmp', '.tiff')):
                contenido = leer_imagen(ruta_archivo)
            elif archivo.endswith('.db'):
                contenido = leer_db(ruta_archivo)
            else:
                errores += 1
                print(f"Tipo de archivo no soportado: {archivo}")
                continue

            if contenido:
                texto_limpio = limpiar_texto(contenido)
                nombre_txt_limpio = os.path.join(carpeta_txt, archivo_limpio)
                with open(nombre_txt_limpio, 'w', encoding='utf-8') as f:
                    f.write(texto_limpio)

                actualizar_log(archivo_limpio, log_file)
                nuevos_procesados += 1
                print(f"Procesado y guardado: {archivo}")

        except Exception as e:
            errores += 1
            print(f"Error procesando archivo {archivo}: {e}")

        processed_size += os.path.getsize(ruta_archivo)
        tqdm.write(f"Progreso: {processed_size / 1024 / 1024:.2f} MB de {total_size / 1024 / 1024:.2f} MB procesados")

    guardar_informe_procesamiento(nuevos_procesados, procesados_saltados, errores, ruta_archivo_informe)

# Función para tokenizar todos los archivos en la carpeta de texto procesado
def tokenizar_todos_archivos(carpeta_txt, carpeta_tokenized):
    for archivo en os.listdir(carpeta_txt):
        if archivo.endswith('_limpio.txt'):
            ruta_archivo = os.path.join(carpeta_txt, archivo)
            nombre_tokens = os.path.join(carpeta_tokenized, archivo.replace('_limpio.txt', '_tokenizado.txt'))
            if not os.path.exists(nombre_tokens):
                with open(ruta_archivo, 'r', encoding='utf-8') as f:
                    contenido = f.read()
                    tokenizar_y_guardar(contenido, nombre_tokens)
    print("Tokenización completada para todos los archivos.")

# Configuraciones de ruta
ruta_base = os.path.dirname(__file__)
ruta_carpeta = os.path.join(ruta_base, 'files')
carpeta_txt = os.path.join(ruta_base, 'txt')
carpeta_tokenized = os.path.join(ruta_base, 'tokenized')
log_file = os.path.join(ruta_base, 'archivos_procesados.log')
ruta_archivo_informe = os.path.join(ruta_base, 'procesado_error.txt')

# Crear las carpetas si no existen
os.makedirs(carpeta_txt, exist_ok=True)
os.makedirs(carpeta_tokenized, exist_ok=True)

# Procesar archivos y tokenizarlos
procesar_archivos(ruta_carpeta, carpeta_txt, log_file, ruta_archivo_informe)
tokenizar_todos_archivos(carpeta_txt, carpeta_tokenized)
